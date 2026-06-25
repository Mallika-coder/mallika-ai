import os
from typing import Dict, Any

from app.tools.base_tool import BaseTool


class DocumentGeneratorTool(BaseTool):
    name = "document_generator"
    description = "Generate documents in various formats: DOCX, PPTX, XLSX, PDF"
    parameters = {
        "type": "object",
        "properties": {
            "format": {
                "type": "string",
                "enum": ["docx", "pptx", "xlsx", "pdf"],
                "description": "Output format",
            },
            "content": {
                "type": "object",
                "description": "Document content structure",
            },
            "filename": {"type": "string", "description": "Output filename"},
        },
        "required": ["format", "content", "filename"],
    }

    async def execute(self, format: str, content: Dict, filename: str) -> Dict[str, Any]:
        output_dir = "/tmp/generated"
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, filename)

        try:
            if format == "docx":
                return await self._generate_docx(content, output_path)
            elif format == "pptx":
                return await self._generate_pptx(content, output_path)
            elif format == "xlsx":
                return await self._generate_xlsx(content, output_path)
            elif format == "pdf":
                return await self._generate_pdf(content, output_path)
            else:
                return {"error": f"Unsupported format: {format}"}
        except Exception as e:
            return {"error": f"Generation failed: {str(e)}"}

    async def _generate_docx(self, content: Dict, output_path: str) -> Dict:
        from docx import Document

        doc = Document()
        for section in content.get("sections", []):
            section_type = section.get("type", "paragraph")
            if section_type == "heading":
                doc.add_heading(section["text"], level=section.get("level", 1))
            elif section_type == "paragraph":
                doc.add_paragraph(section["text"])
            elif section_type == "table":
                data = section["data"]
                if data:
                    table = doc.add_table(rows=len(data), cols=len(data[0]))
                    for i, row in enumerate(data):
                        for j, cell in enumerate(row):
                            table.rows[i].cells[j].text = str(cell)
            elif section_type == "list":
                for item in section.get("items", []):
                    doc.add_paragraph(item, style="List Bullet")
        doc.save(output_path)
        return {"path": output_path, "filename": os.path.basename(output_path), "format": "docx"}

    async def _generate_pptx(self, content: Dict, output_path: str) -> Dict:
        from pptx import Presentation

        prs = Presentation()
        for slide_data in content.get("slides", []):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            if slide_data.get("content"):
                slide.placeholders[1].text = slide_data["content"]
        prs.save(output_path)
        return {"path": output_path, "filename": os.path.basename(output_path), "format": "pptx"}

    async def _generate_xlsx(self, content: Dict, output_path: str) -> Dict:
        from openpyxl import Workbook

        wb = Workbook()
        first = True
        for sheet_data in content.get("sheets", []):
            if first:
                ws = wb.active
                ws.title = sheet_data.get("name", "Sheet1")
                first = False
            else:
                ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))
            for row in sheet_data.get("data", []):
                ws.append(row)
        wb.save(output_path)
        return {"path": output_path, "filename": os.path.basename(output_path), "format": "xlsx"}

    async def _generate_pdf(self, content: Dict, output_path: str) -> Dict:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas as pdf_canvas
        from reportlab.lib.units import inch

        c = pdf_canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        y = height - inch

        for section in content.get("sections", []):
            if y < inch:
                c.showPage()
                y = height - inch

            section_type = section.get("type", "paragraph")
            if section_type == "heading":
                c.setFont("Helvetica-Bold", 16)
                c.drawString(inch, y, section["text"])
                y -= 30
            elif section_type == "paragraph":
                c.setFont("Helvetica", 11)
                words = section["text"].split()
                line = ""
                for word in words:
                    if c.stringWidth(line + " " + word, "Helvetica", 11) < width - 2 * inch:
                        line += " " + word if line else word
                    else:
                        c.drawString(inch, y, line)
                        y -= 15
                        line = word
                if line:
                    c.drawString(inch, y, line)
                    y -= 20

        c.save()
        return {"path": output_path, "filename": os.path.basename(output_path), "format": "pdf"}
